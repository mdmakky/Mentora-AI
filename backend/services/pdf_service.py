import fitz  # PyMuPDF
import hashlib
import io
import logging
from typing import List, Tuple, Optional

logger = logging.getLogger(__name__)

# Minimum characters on a page to consider it "has real text".
# Pages with fewer (e.g. just a page number or header) are treated as image pages.
_MIN_TEXT_CHARS = 20

# Set to True the first time Tesseract is found to be unavailable,
# so we stop attempting OCR for the rest of the process lifetime.
_tesseract_unavailable: bool = False


def extract_text_from_pdf(file_bytes: bytes) -> List[dict]:
    """Extract text from PDF page by page (no OCR)."""
    pages = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text("text")
        if text.strip():
            pages.append({
                "page_number": page_num + 1,
                "content": text.strip(),
            })

    doc.close()
    return pages


def extract_text_from_pdf_with_ocr(file_bytes: bytes) -> Tuple[List[dict], bool]:
    """
    Extract text from a PDF with per-page Tesseract OCR fallback.

    For every page:
      - First attempt: standard text layer extraction (PyMuPDF).
      - If the page has fewer than _MIN_TEXT_CHARS characters (image-only or
        near-blank), it is re-processed via PyMuPDF's built-in Tesseract OCR
        bridge (page.get_textpage_ocr).

    Returns:
        (pages, ocr_applied)
        - pages: same List[{"page_number": int, "content": str}] format as
          extract_text_from_pdf, ready for chunk_document().
        - ocr_applied: True if OCR was used on at least one page.

    Requires: tesseract-ocr installed on the system (`apt install tesseract-ocr`).
    If Tesseract is not present, image pages are silently skipped and
    ocr_applied is always False (no error raised).
    """
    global _tesseract_unavailable

    pages: List[dict] = []
    ocr_applied = False

    doc = fitz.open(stream=file_bytes, filetype="pdf")

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text("text").strip()

        if len(text) >= _MIN_TEXT_CHARS:
            # Normal text layer — use as-is
            pages.append({
                "page_number": page_num + 1,
                "content": text,
            })
            continue

        # Image-only / sparse page — attempt OCR
        if _tesseract_unavailable:
            # Tesseract confirmed missing earlier; skip silently
            if text:
                pages.append({"page_number": page_num + 1, "content": text})
            continue

        try:
            ocr_tp = page.get_textpage_ocr(
                flags=3,       # preserve whitespace + ligatures
                language="eng",
                dpi=150,       # 150 DPI matches quality used elsewhere in the pipeline
                full=True,     # OCR the full page (not just detected image regions)
            )
            ocr_text = page.get_text(textpage=ocr_tp).strip()

            if ocr_text and len(ocr_text) >= _MIN_TEXT_CHARS:
                pages.append({
                    "page_number": page_num + 1,
                    "content": ocr_text,
                })
                ocr_applied = True
                logger.debug(
                    "[pdf_service] OCR applied to page %d (%d chars extracted)",
                    page_num + 1, len(ocr_text),
                )
            elif text:
                # OCR yielded nothing useful, at least keep the sparse original text
                pages.append({"page_number": page_num + 1, "content": text})

        except RuntimeError as e:
            err_lower = str(e).lower()
            if "tesseract" in err_lower or "tessdata" in err_lower or "no ocr support" in err_lower:
                _tesseract_unavailable = True
                logger.warning(
                    "[pdf_service] Tesseract OCR is unavailable (%s) — "
                    "image-only pages will be skipped. "
                    "Install with: sudo apt install tesseract-ocr",
                    e,
                )
            else:
                logger.warning("[pdf_service] OCR failed on page %d: %s", page_num + 1, e)
            if text:
                pages.append({"page_number": page_num + 1, "content": text})

        except Exception as e:
            logger.warning("[pdf_service] OCR error on page %d: %s", page_num + 1, e)
            if text:
                pages.append({"page_number": page_num + 1, "content": text})

    doc.close()
    return pages, ocr_applied


def get_pdf_page_count(file_bytes: bytes) -> int:
    """Get the number of pages in a PDF."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    count = len(doc)
    doc.close()
    return count


def get_pdf_metadata(file_bytes: bytes) -> dict:
    """Extract PDF metadata for copyright check."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    metadata = doc.metadata or {}
    doc.close()
    return metadata


def extract_text_from_docx(file_bytes: bytes) -> List[dict]:
    """Extract text from DOCX file."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        pages = []
        current_text = []
        page_num = 1

        for para in doc.paragraphs:
            current_text.append(para.text)
            # Approximate page breaks every ~3000 chars
            if len("\n".join(current_text)) > 3000:
                pages.append({
                    "page_number": page_num,
                    "content": "\n".join(current_text).strip(),
                })
                current_text = []
                page_num += 1

        if current_text:
            pages.append({
                "page_number": page_num,
                "content": "\n".join(current_text).strip(),
            })

        return pages
    except Exception as e:
        return [{"page_number": 1, "content": f"Error extracting DOCX: {str(e)}"}]


def extract_text_from_pptx(file_bytes: bytes) -> List[dict]:
    """Extract text from PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []

        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                slides.append({
                    "page_number": slide_num,
                    "content": "\n".join(texts).strip(),
                    "slide_number": slide_num,
                })

        return slides
    except Exception as e:
        return [{"page_number": 1, "content": f"Error extracting PPTX: {str(e)}"}]


def extract_text_from_image(file_bytes: bytes, file_ext: str) -> Tuple[List[dict], bool]:
    """
    Run OCR on a standalone JPEG or PNG image using PyMuPDF + Tesseract.

    The image is wrapped into a one-page in-memory PDF so the existing
    OCR-aware PDF path (`extract_text_from_pdf_with_ocr`) can be reused
    without duplicating any logic.

    Returns:
        (pages, ocr_applied) — same contract as extract_text_from_pdf_with_ocr.
    """
    ext = file_ext.lower().lstrip('.')
    fitz_type = "jpeg" if ext in ("jpg", "jpeg") else "png"
    try:
        img_doc = fitz.open(stream=file_bytes, filetype=fitz_type)
        pdf_bytes = img_doc.convert_to_pdf()
        img_doc.close()
    except Exception as e:
        logger.warning("[pdf_service] Could not open image '%s' for OCR: %s", file_ext, e)
        return [], False

    return extract_text_from_pdf_with_ocr(pdf_bytes)


def calculate_file_hash(file_bytes: bytes) -> str:
    """Calculate SHA-256 hash of file content."""
    return hashlib.sha256(file_bytes).hexdigest()
